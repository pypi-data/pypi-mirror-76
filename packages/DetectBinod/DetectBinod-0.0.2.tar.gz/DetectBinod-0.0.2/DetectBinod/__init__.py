from PyPDF2 import PdfFileReader
import docx 
import xlrd
from pptx import Presentation

def status(p,filename):
    # Printing status
    if(p==1):
        print(f"Binod Found in file {filename}")
    else:
        print(f"No Binod in {filename} bro!\nCongrats !")
    return 
    

def isBinod_xlsx(filename):

    '''
    To find binod in xlsx file
    arg: filename as string i.e. path it will be opened in rb mode to read
    return nothing just updates status i.e. prints on terminal about बीनोडता  
    '''
    book=xlrd.open_workbook(filename)
    sheet=book.sheet_by_index(0)
    i=sheet.nrows
    j=sheet.ncols
    # print(i,j)
    for p in range(i):
        for q in range(j):
            if(sheet.cell_value(p,q)=="binod"):
                status(1,filename)
                print(f"First Binod is found at {(p+1),(q+1)} cell")
                return 
            else:
                pass
            status(0,filename)
    return


def isBinod_txt(filename):
    '''
    Function to detect Binod in txt file-format
    arg: file-name or path of file, opened in rb mode
    calls status and prints on terminal 
    '''

    with open(filename) as f:
        fileContent=f.read()
        if "binod" in fileContent.lower():
            status(1,filename)
        else:
            status(0,filename)
    return        



def isBinod_pdf(filename):
    ''' 
    Function to detect Binod in pdf and call status according to it
    '''
    f= open(filename,"rb")
    pdf=PdfFileReader(f)
    for i in range(pdf.numPages):
        page=pdf.getPage(i)
        cont=str(page.extractText()).lower()
        if "binod" in cont:
            status(1,filename)
        else:
            status(0,filename)
    return



def isBinod_docx(filename):
    doc=docx.Document(filename)
    k=len(doc.paragraphs)

    cnt=0
    for i in range(k):
        for j in range(len(doc.paragraphs[i].runs)):
            if str(doc.paragraphs[i].runs[j].text).lower()=="binod":
                status(1,filename)
                return 
    
    status(0,filename)            
    return 
    


def isBinod_ppt(filename):
    ppt = Presentation(filename)
    File_to_write_data= open("File_To_Extract_ppt.txt","w")
        
    for slide in ppt.slides: 
        for shape in slide.shapes: 
            if not shape.has_text_frame: 
                continue 
            for paragraph in shape.text_frame.paragraphs: 
                for run in paragraph.runs: 
                    if "binod" in run.text.lower():
                       status(1,filename)
                       return 
                    else:
                        status(0,filename)
                        return 

    # isBinod_txt(File_to_write_data)
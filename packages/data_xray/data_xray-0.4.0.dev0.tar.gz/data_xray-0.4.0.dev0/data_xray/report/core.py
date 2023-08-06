from ..scan import PlotImage
from ..grid import *
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from matplotlib import pyplot as plt
import pandas as pd

#some utilies to work with powerpoint
meanstd = lambda arr: np.std(np.ravel(arr)/np.mean(np.ravel(arr)))

class SummaryPPT(object):
    def __init__(self, pname="image_summary", new=False, fdict=None, maximages=50, chanselect = {'scan':'Z', 'grid':'cf'}, **kwargs):

        if fdict is None:
            print('please specify data to summarize')
            return
        else:
            fdicts = [fdict[i:i + maximages] for i in range(0, len(fdict), maximages)]

        self.topdir = os.path.commonpath([j.fname for j in fdict]) + '/'
        self.chanselect = chanselect
        for j,f in enumerate(fdicts): 
            self.presentation_name = pname + '_' + str(j)
            self.pptx_file_name = self.topdir + self.presentation_name + '.pptx'
            self.fdict = f
            #self.new = new
            self.init_ppt(self.presentation_name)
            
            #self.insert_images()
            self.insert_data()
            #try:
            self.pres.save(self.pptx_file_name)
            print('batch ' + str(j) +' stored in : ' + self.pptx_file_name )
            #except:
            #    print('something wrong with saving the presentation file ' + self.pptx_file_name)

    def init_ppt(self, presentation_name):
        pres = Presentation()
        pres.notes_master.name = self.presentation_name
        self.pres = pres

    def insert_data(self):
        for fj in self.fdict:
            if re.findall('sxm', fj.fname):
                try:
                    if self.chanselect == "Automatic":
                        # attempt to recognize good data
                        plotsignals = []
                        for c in fj.signals.keys():
                            sig = fj.signals[c]['forward']
                            sig = sig[~np.isnan(sig)]
                            if meanstd(sig) > 2:
                                plotsignals.append(c)
                        plotsignals.append('Z')
                        plotsignals = list(set(plotsignals))  # no dobule Z

                    else:
                        plotsignals = self.chanselect['scan']

                    nrows = 2 if len(plotsignals) > 2 else 1
                    ncols = int(np.ceil(len(plotsignals) / nrows))
                    f3, a3 = plt.subplots(nrows, ncols);
                    if a3 is not list:
                        a3 = [a3]
                    for c, a in zip(plotsignals, np.ravel(a3)):
                        PlotImage(fj, chan=c, ax=a, high_pass=None);

                    [a.axis('off') for a in np.ravel(a3)]

                    xy = ['X', 'Y']
                    offset = fj.header['scan_offset'] / 1e-9
                    xyoffsets = [xy[j] + '=' + str(np.round(offset[j], 2)) + ' nm ' for j in [0, 1]]

                    titleString = [fj.fname]
                    titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
                    titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
                    titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
                    titleString.append('Resolution: ' + str(fj.header['scan_pixels']))

                    self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                    print(os.path.basename(fj.ds.fname) + ' imported')
                    f3.clf();  # close figure so that it doesn't clog up in the end
                except:
                    print(os.path.basename(fj.fname) + ' failed')

            elif re.findall('3ds', fj.fname):
                try:
                    fig, ax = plt.subplots(1, 2)
                    ChanHistogramDS(fj.ds, xy=['bias', self.chanselect['grid']], xymod=[lambda x: x, lambda x: x / np.mean(np.ravel(x))],
                                    ax=ax[0], label=['bias', self.chanselect['grid'], ''])

                    # plop in a clustered map
                    km = ChanPcaKmeansDS(fj.ds, xvec='bias', chan=self.chanselect['grid'], mod=lambda x: x / np.mean(np.ravel(x)),
                                         comps=6, nclust=4, fig=None)
                    ax[1].imshow(km)
                    #plt.colorbar()

                    titleString = [fj.fname]
                except:
                    print(os.path.basename(fj.fname) + ' failed. Import skipped')
                    continue
                # titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
                # titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
                # titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
                # titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
                try:
                    self.fig_to_ppt([fig], leftop=[1, 2], txt=titleString)
                    print(os.path.basename(fj.fname) + ' imported')

                    fig.clf();  # close figure so that it doesn't clog up in the end

                except:
                    print(os.path.basename(fj.fname) + ' import into ppt failed')

    def insert_maps(self, mapchan=['cf']):
        
        #plop in a histogram of the map channel across the whole thing
        
        for fj in self.fdict:
            fig,ax=plt.subplots(2,1, figsize=(8,5))
            ChanHistogramDS(fj.ds, xy=['bias',mapchan[0]], xymod=[lambda x:x,lambda x:x/np.mean(np.ravel(x))], ax=ax[0], label=['bias',mapchan[0][:-1],''])  
            
            #plop in a clustered map
            km = ChanPcaKmeansDS(fj.ds, xvec='bias', chan=mapchan[0], mod = lambda x: x/np.mean(np.ravel(x)), comps=6, nclust=4, fig=None)
            ax[1].imshow(km)
            plt.colorbar()
        
            #fig.savefig("pca_3d.png",bbox_inches='tight')

            ###need to add name attribute to grids
            titleString = [fj.fname]
        

            # titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
            # titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
            # titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
            # titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
            try:
                self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                print(os.path.basename(fj.ds.fname) + ' imported')
            
                f3.clf(); #close figure so that it doesn't clog up in the end
            
            except:
                print(os.path.basename(fj.ds.fname) + ' failed')
            
            #self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)

        return

    def insert_images(self):
        """
        Dumpt a batch of images into a powerpoint

        :param chanselect:
        :param fdict:
        :param topdir:
        :return:
        """
        #for folder in (self.fdict.keys()):
        #        self.text_to_slide(folder)

        # newpres = Presentation()
        # newpres.notes_master.name = 'sum1.pptx'
        # newpres.save(newpres.notes_master.name)

        for fj in self.fdict:
            # TextToSlide(fj.fname,pres=pres)
            try:
                # scf = {1:(12,6), 2:(12,9), 3:(12,10)}
                if self.chanselect == "Automatic":
                #attempt to recognize good data  
                    plotsignals = []
                    for c in fj.signals.keys():
                        sig = fj.signals[c]['forward']
                        sig = sig[~np.isnan(sig)]
                        if meanstd(sig) > 2:
                            plotsignals.append(c)    
                    plotsignals.append('Z')
                    plotsignals = list(set(plotsignals)) #no dobule Z
                    

                else:
                    plotsignals = self.chanselect
                    
                nrows = 2 if len(plotsignals) > 2 else 1
                ncols = int(np.ceil(len(plotsignals)/nrows))
                f3, a3 = plt.subplots(nrows,ncols);
                if a3 is not list:
                    a3 = [a3]
                for c,a in zip(plotsignals,np.ravel(a3)):
                    PlotImage(fj, chan=c, ax=a, high_pass=None);
                
                [a.axis('off') for a in np.ravel(a3)]

               

                xy = ['X','Y']
                offset = fj.header['scan_offset']/1e-9
                xyoffsets = [xy[j] + '=' + str(np.round(offset[j],2)) + ' nm ' for j in [0,1]]

                titleString = [fj.fname]
                titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
                titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
                titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
                titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
                
                self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                print(os.path.basename(fj.ds.fname) + ' imported')
                f3.clf(); #close figure so that it doesn't clog up in the end
            except:
                print(os.path.basename(fj.ds.fname) + ' failed')

    def png_to_ppt(self, pngfile, ttl = []):
       """
       Plop a PNG file into powerpoint slide
       :param pngfile:
       :param pres:
       :param ttl:
       :return:
       """

       #blank_slide_layout = pres.slide_layouts[6]
       title_slide_layout = self.pres.slide_layouts[9]

       left = top = Inches(1)

       slide = self.pres.slides.add_slide(title_slide_layout)
       slide.shapes.add_picture(pngfile, left, top)
       subtitle = slide.placeholders[1]
       title = slide.shapes.title
       if len(ttl):
           subtitle.text = ttl


    def fig_to_ppt(self, figs, leftop=[0,1.5], txt=None):
        """
        Plop figures into powerpoint
        :param figs:
        :param pres:
        :param leftop:
        :param txt:
        :return:
        """
        #savepptx needs to be a full path. If None is provided the default presentation
        #will be created with a name sum1.pptx in the current folder
        from pptx.util import Inches

        blank_slide_layout = self.pres.slide_layouts[5]
        left = Inches(leftop[0])
        top = Inches(leftop[1])

        tmp_path = 't1.png'
        for figp in figs:
            plt.savefig(tmp_path, transparent=1, format='png', dpi=300, bbox_inches = 'tight')
            slide = self.pres.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(tmp_path, left, top)

        if txt is not None:
            self.text_to_slide(txt, slide=slide)

    def text_to_slide(self, txt, slide=None): #lets make txt a list of strings
                                              
        """
        convert text to slide

        :param txt: list of strings
        :param pres:
        :param slide:
        :return:
        """
        from pptx.util import Pt
        #title = slide.shapes.title
        #subtitle = slide.placeholders[1]

       # title.text = "Hello, World!"
        #subtitle.text = "python-pptx was here!"

       # prs.save('test.pptx')

        from pptx.util import Inches

        if self.pres == None:
            print('please init presentation')
        else:
            if slide is None:
                bullet_slide_layout = self.pres.slide_layouts[5]
                slide = self.pres.slides.add_slide(bullet_slide_layout)

            shapes = slide.shapes

            countshapes = 0

            #just catch the first shape object with a frame in it
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elif countshapes > 0:
                    tframe = shape.text_frame
                    tframe.clear()
                    #print('caught one')
                else:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    countshapes = 1

            text_frame.clear()
            p = text_frame.paragraphs[0]
            
            for t in txt:
                run = p.add_run()
                run.text = t

                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)    
                #    p.text = t
                p = text_frame.add_paragraph()

                 
       
            


# class SummaryPPT2(object):
#     def __init__(self, pname='summary', new=False, fdict=None, topdir=None, **kwargs):

#         if fdict is None:
#             print('please specify data to summarize')
#             return
#         else:
#             self.fdict = fdict
#             self.pname = pname
#             self.new = new
#             if topdir==None:
#                 topdir = os.getcwd()


#     def crawl_and_save(self):
#         for ind, (folder, files) in enumerate(zip(fdict.keys(), fdict.items())):
            
#             #self.topdir = os.path.commonpath([j.ds.fname for j in fdict]) + '/'
#             self.presentation_name = pname
#             self.pptx_file_name = folder + self.presentation_name + '.pptx'

#             self.new = new
#             self.init_ppt()
#             self.insert_images()

#             try:
#                 self.pres.save(self.pptx.file_name)
#                 print('images stored in : ' + self.pptx_file_name )
#             except:
#                 print('something wrong with saving the presentation file')


#     def init_ppt(self):
#         """
#         Initialize powerpoint presentation
#         :param pname:
#         :param new:
#         :return:
#         """
#         if self.presentation_name is None:
#             newpres = Presentation()
#             newpres.notes_master.name = 'sum1.pptx'
#             newpres.save(newpres.notes_master.name)
#             self.pres = newpres
#         elif self.new:
#             newpres = Presentation()
#             newpres.notes_master.name = self.presentation_name
#             newpres.save(newpres.notes_master.name)
#             self.pres = newpres
#         else:
#             pres = Presentation(self.presentation_name)
#             pres.notes_master.name = self.presentation_name
#             self.pres = pres

#     def insert_images(self, chanselect='Z'):
#         """
#         Dumpt a batch of images into a powerpoint

#         :param chanselect:
#         :param fdict:
#         :param topdir:
#         :return:
#         """
#         #for folder in (self.fdict.keys()):
#         #        self.text_to_slide(folder)

#         for fj in self.fdict:
#             # TextToSlide(fj.fname,pres=pres)
#             try:
#                 f3, a3 = plt.subplots(1, 1)
#                 d2 = PlotImage(fj, chan=chanselect, ax=a3, high_pass=None)
#                 self.fig_to_ppt([f3], leftop=[3, 2], txt=fj.fname)
#                 print(os.path.basename(fj.ds.fname) + ' imported')

#             except:
#                 print(os.path.basename(fj.ds.fname) + ' failed')

#     def png_to_ppt(self, pngfile, ttl = []):
#        """
#        Plop a PNG file into powerpoint slide
#        :param pngfile:
#        :param pres:
#        :param ttl:
#        :return:
#        """

#        #blank_slide_layout = pres.slide_layouts[6]
#        title_slide_layout = self.pres.slide_layouts[9]

#        left = top = Inches(1)

#        slide = self.pres.slides.add_slide(title_slide_layout)
#        slide.shapes.add_picture(pngfile, left, top)
#        subtitle = slide.placeholders[1]
#        title = slide.shapes.title
#        if len(ttl):
#            subtitle.text = ttl


#     def fig_to_ppt(self, figs, leftop=[0,1.5], txt=None):
#         """
#         Plop figures into powerpoint
#         :param figs:
#         :param pres:
#         :param leftop:
#         :param txt:
#         :return:
#         """
#         #savepptx needs to be a full path. If None is provided the default presentation
#         #will be created with a name sum1.pptx in the current folder
#         from pptx.util import Inches

#         blank_slide_layout = self.pres.slide_layouts[5]
#         left = Inches(leftop[0])
#         top = Inches(leftop[1])

#         tmp_path = 't1.png'
#         for figp in figs:
#             plt.savefig(tmp_path, transparent=1, format='png', dpi=300, bbox_inches = 'tight')
#             slide = self.pres.slides.add_slide(blank_slide_layout)
#             slide.shapes.add_picture(tmp_path, left, top)

#         if txt is not None:
#             self.text_to_slide(txt, slide=slide)

#     def text_to_slide(self, txt, slide=None):
#         """
#         convert text to slide

#         :param txt:
#         :param pres:
#         :param slide:
#         :return:
#         """
#         from pptx.util import Pt
#         #title = slide.shapes.title
#         #subtitle = slide.placeholders[1]

#        # title.text = "Hello, World!"
#         #subtitle.text = "python-pptx was here!"

#        # prs.save('test.pptx')

#         from pptx.util import Inches

#         if self.pres == None:
#             print('please init presentation')
#         else:
#             if slide is None:
#                 bullet_slide_layout = self.pres.slide_layouts[5]
#                 slide = self.pres.slides.add_slide(bullet_slide_layout)

#             shapes = slide.shapes

#             countshapes = 0

#             #just catch the first shape object with a frame in it
#             for shape in slide.shapes:
#                 if not shape.has_text_frame:
#                     continue
#                 elif countshapes > 0:
#                     tframe = shape.text_frame
#                     tframe.clear()
#                     #print('caught one')
#                 else:
#                     text_frame = shape.text_frame
#                     text_frame.clear()
#                     countshapes = 1

#             text_frame.clear()
#             p = text_frame.paragraphs[0]
#             run = p.add_run()
#             run.text = txt

#             font = run.font
#             font.name = 'Calibri'
#             font.size = Pt(12)


